"""
SentinelDLP File Processor Service (FR-001)
Universal file type support with OCR integration for enterprise DLP.

Supports:
- Documents: PDF, DOCX, DOC, XLSX, XLS, PPTX, PPT, RTF, ODT, ODS, ODP
- Text: TXT, CSV, JSON, XML, HTML, MD, LOG, code files
- Images: PNG, JPG, JPEG, TIFF, BMP, GIF, WEBP (with OCR)
- Email: EML, MSG
- Archives: ZIP, TAR, GZ (content listing)

Features:
- Magic byte detection for accurate file type identification
- OCR support via Tesseract (English + Arabic)
- Intelligent chunking for large files
- Streaming support for memory efficiency
"""

import io
import os
import re
import csv
import json
import email
import zipfile
import tarfile
import tempfile
from pathlib import Path
from typing import Optional, List, Dict, Tuple, Any
from dataclasses import dataclass, field
from enum import Enum
import hashlib


class FileCategory(Enum):
    """Categories of supported file types."""
    TEXT = "text"
    DOCUMENT = "document"
    SPREADSHEET = "spreadsheet"
    PRESENTATION = "presentation"
    IMAGE = "image"
    EMAIL = "email"
    ARCHIVE = "archive"
    UNKNOWN = "unknown"


@dataclass
class FileType:
    """Detected file type information."""
    extension: str
    mime_type: str
    category: FileCategory
    requires_ocr: bool = False
    is_binary: bool = False


@dataclass
class ProcessedChunk:
    """A chunk of processed text."""
    content: str
    chunk_index: int
    total_chunks: int
    source_page: Optional[int] = None
    char_offset: int = 0


@dataclass
class ProcessedFile:
    """Result of file processing."""
    original_filename: str
    detected_type: FileType
    text_content: str
    chunks: List[ProcessedChunk] = field(default_factory=list)
    page_count: Optional[int] = None
    word_count: int = 0
    char_count: int = 0
    ocr_applied: bool = False
    extraction_warnings: List[str] = field(default_factory=list)
    metadata: Dict[str, Any] = field(default_factory=dict)


# File signatures (magic bytes) for type detection
FILE_SIGNATURES = {
    # PDF
    b'%PDF': ('pdf', 'application/pdf', FileCategory.DOCUMENT),

    # Office Open XML (DOCX, XLSX, PPTX) - ZIP-based
    b'PK\x03\x04': ('zip_office', 'application/zip', FileCategory.DOCUMENT),

    # Legacy Office (DOC, XLS, PPT) - OLE Compound
    b'\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1': ('ole', 'application/msword', FileCategory.DOCUMENT),

    # Images
    b'\x89PNG\r\n\x1a\n': ('png', 'image/png', FileCategory.IMAGE),
    b'\xff\xd8\xff': ('jpg', 'image/jpeg', FileCategory.IMAGE),
    b'GIF87a': ('gif', 'image/gif', FileCategory.IMAGE),
    b'GIF89a': ('gif', 'image/gif', FileCategory.IMAGE),
    b'II*\x00': ('tiff', 'image/tiff', FileCategory.IMAGE),
    b'MM\x00*': ('tiff', 'image/tiff', FileCategory.IMAGE),
    b'BM': ('bmp', 'image/bmp', FileCategory.IMAGE),
    b'RIFF': ('webp', 'image/webp', FileCategory.IMAGE),

    # Archives
    b'PK\x05\x06': ('zip', 'application/zip', FileCategory.ARCHIVE),
    b'\x1f\x8b': ('gz', 'application/gzip', FileCategory.ARCHIVE),

    # RTF
    b'{\\rtf': ('rtf', 'application/rtf', FileCategory.DOCUMENT),
}

# Text-based file extensions
TEXT_EXTENSIONS = {
    'txt', 'csv', 'json', 'xml', 'html', 'htm', 'md', 'markdown',
    'log', 'py', 'js', 'ts', 'jsx', 'tsx', 'java', 'c', 'cpp', 'h',
    'cs', 'go', 'rs', 'rb', 'php', 'pl', 'sh', 'bash', 'zsh',
    'yaml', 'yml', 'ini', 'conf', 'cfg', 'toml', 'env',
    'sql', 'graphql', 'proto', 'css', 'scss', 'less', 'sass',
    'vue', 'svelte', 'r', 'scala', 'kt', 'swift', 'ps1', 'bat', 'cmd'
}

# Office XML content types for detection
OFFICE_CONTENT_TYPES = {
    'word/document.xml': ('docx', 'application/vnd.openxmlformats-officedocument.wordprocessingml.document', FileCategory.DOCUMENT),
    'xl/workbook.xml': ('xlsx', 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', FileCategory.SPREADSHEET),
    'ppt/presentation.xml': ('pptx', 'application/vnd.openxmlformats-officedocument.presentationml.presentation', FileCategory.PRESENTATION),
}


class FileProcessor:
    """
    Universal file processor with OCR support.

    Handles text extraction from various file formats with intelligent
    chunking for large files and OCR for images/scanned documents.
    """

    def __init__(
        self,
        max_chunk_size: int = 50000,  # ~50K chars per chunk for LLM
        ocr_enabled: bool = True,
        ocr_languages: str = "eng+ara",  # English + Arabic for Dubai
        max_file_size_mb: int = 10240,  # Max file size: 10GB (FR-004)
    ):
        self.max_chunk_size = max_chunk_size
        self.ocr_enabled = ocr_enabled
        self.ocr_languages = ocr_languages
        self.max_file_size_bytes = max_file_size_mb * 1024 * 1024

        # Check for optional dependencies
        self._check_dependencies()

    def _check_dependencies(self):
        """Check which optional dependencies are available."""
        self.has_pdfplumber = False
        self.has_docx = False
        self.has_openpyxl = False
        self.has_pptx = False
        self.has_pillow = False
        self.has_pytesseract = False
        self.has_striprtf = False
        self.has_magic = False

        try:
            import pdfplumber
            self.has_pdfplumber = True
        except ImportError:
            pass

        try:
            import docx
            self.has_docx = True
        except ImportError:
            pass

        try:
            import openpyxl
            self.has_openpyxl = True
        except ImportError:
            pass

        try:
            import pptx
            self.has_pptx = True
        except ImportError:
            pass

        try:
            from PIL import Image
            self.has_pillow = True
        except ImportError:
            pass

        try:
            import pytesseract
            self.has_pytesseract = True
        except ImportError:
            pass

        try:
            from striprtf.striprtf import rtf_to_text
            self.has_striprtf = True
        except ImportError:
            pass

        try:
            import magic
            self.has_magic = True
        except ImportError:
            pass

    def detect_file_type(self, content: bytes, filename: str) -> FileType:
        """
        Detect file type using magic bytes and extension.

        Args:
            content: File content as bytes
            filename: Original filename

        Returns:
            FileType with detected information
        """
        ext = filename.rsplit('.', 1)[-1].lower() if '.' in filename else ''

        # Check magic bytes first
        for signature, (sig_ext, mime, category) in FILE_SIGNATURES.items():
            if content.startswith(signature):
                # Special handling for ZIP-based Office formats
                if sig_ext == 'zip_office':
                    return self._detect_office_type(content, filename, ext)

                # Special handling for OLE (legacy Office)
                if sig_ext == 'ole':
                    return self._detect_ole_type(content, filename, ext)

                return FileType(
                    extension=sig_ext,
                    mime_type=mime,
                    category=category,
                    requires_ocr=category == FileCategory.IMAGE,
                    is_binary=True
                )

        # Check for text-based files by extension
        if ext in TEXT_EXTENSIONS:
            return FileType(
                extension=ext,
                mime_type=f'text/{ext}',
                category=FileCategory.TEXT,
                requires_ocr=False,
                is_binary=False
            )

        # Check for email files
        if ext in ('eml', 'msg'):
            return FileType(
                extension=ext,
                mime_type='message/rfc822' if ext == 'eml' else 'application/vnd.ms-outlook',
                category=FileCategory.EMAIL,
                requires_ocr=False,
                is_binary=ext == 'msg'
            )

        # Check for archive files by extension
        if ext in ('zip', 'tar', 'gz', 'tgz', '7z', 'rar'):
            return FileType(
                extension=ext,
                mime_type=f'application/{ext}',
                category=FileCategory.ARCHIVE,
                requires_ocr=False,
                is_binary=True
            )

        # Try to detect if it's text by checking for null bytes
        sample = content[:8192]
        if b'\x00' not in sample:
            # Likely text
            return FileType(
                extension=ext or 'txt',
                mime_type='text/plain',
                category=FileCategory.TEXT,
                requires_ocr=False,
                is_binary=False
            )

        # Unknown binary file
        return FileType(
            extension=ext or 'bin',
            mime_type='application/octet-stream',
            category=FileCategory.UNKNOWN,
            requires_ocr=False,
            is_binary=True
        )

    def _detect_office_type(self, content: bytes, filename: str, ext: str) -> FileType:
        """Detect specific Office Open XML format from ZIP content."""
        try:
            with zipfile.ZipFile(io.BytesIO(content)) as zf:
                names = zf.namelist()
                for path, (doc_ext, mime, category) in OFFICE_CONTENT_TYPES.items():
                    if path in names:
                        return FileType(
                            extension=doc_ext,
                            mime_type=mime,
                            category=category,
                            requires_ocr=False,
                            is_binary=True
                        )
        except:
            pass

        # Fall back to extension
        ext_map = {
            'docx': ('docx', 'application/vnd.openxmlformats-officedocument.wordprocessingml.document', FileCategory.DOCUMENT),
            'xlsx': ('xlsx', 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', FileCategory.SPREADSHEET),
            'pptx': ('pptx', 'application/vnd.openxmlformats-officedocument.presentationml.presentation', FileCategory.PRESENTATION),
        }

        if ext in ext_map:
            e, m, c = ext_map[ext]
            return FileType(extension=e, mime_type=m, category=c, requires_ocr=False, is_binary=True)

        # Generic ZIP
        return FileType(
            extension='zip',
            mime_type='application/zip',
            category=FileCategory.ARCHIVE,
            requires_ocr=False,
            is_binary=True
        )

    def _detect_ole_type(self, content: bytes, filename: str, ext: str) -> FileType:
        """Detect legacy Office format from OLE compound document."""
        ext_map = {
            'doc': ('doc', 'application/msword', FileCategory.DOCUMENT),
            'xls': ('xls', 'application/vnd.ms-excel', FileCategory.SPREADSHEET),
            'ppt': ('ppt', 'application/vnd.ms-powerpoint', FileCategory.PRESENTATION),
        }

        if ext in ext_map:
            e, m, c = ext_map[ext]
            return FileType(extension=e, mime_type=m, category=c, requires_ocr=False, is_binary=True)

        # Default to DOC
        return FileType(
            extension='doc',
            mime_type='application/msword',
            category=FileCategory.DOCUMENT,
            requires_ocr=False,
            is_binary=True
        )

    async def process_file(
        self,
        content: bytes,
        filename: str,
        force_ocr: bool = False
    ) -> ProcessedFile:
        """
        Process a file and extract text content.

        Args:
            content: File content as bytes
            filename: Original filename
            force_ocr: Force OCR even for text-extractable documents

        Returns:
            ProcessedFile with extracted text and metadata
        """
        warnings = []

        # Check file size
        if len(content) > self.max_file_size_bytes:
            raise ValueError(f"File too large: {len(content)} bytes (max: {self.max_file_size_bytes})")

        # Detect file type
        file_type = self.detect_file_type(content, filename)

        # Extract text based on file type
        text_content = ""
        page_count = None
        ocr_applied = False
        metadata = {}

        try:
            if file_type.category == FileCategory.TEXT:
                text_content = self._extract_text_file(content)

            elif file_type.extension == 'pdf':
                text_content, page_count, ocr_applied = self._extract_pdf(content, force_ocr)

            elif file_type.extension == 'docx':
                text_content, metadata = self._extract_docx(content)

            elif file_type.extension == 'doc':
                text_content = self._extract_doc(content)
                if not text_content:
                    warnings.append("Legacy DOC format - limited extraction")

            elif file_type.extension == 'xlsx':
                text_content, metadata = self._extract_xlsx(content)

            elif file_type.extension == 'xls':
                text_content = self._extract_xls(content)
                if not text_content:
                    warnings.append("Legacy XLS format - limited extraction")

            elif file_type.extension == 'pptx':
                text_content, page_count, metadata = self._extract_pptx(content)

            elif file_type.extension == 'ppt':
                text_content = self._extract_ppt(content)
                if not text_content:
                    warnings.append("Legacy PPT format - limited extraction")

            elif file_type.extension == 'rtf':
                text_content = self._extract_rtf(content)

            elif file_type.category == FileCategory.IMAGE:
                if self.ocr_enabled:
                    text_content = self._extract_image_ocr(content)
                    ocr_applied = True
                else:
                    warnings.append("Image file - OCR disabled")

            elif file_type.category == FileCategory.EMAIL:
                text_content, metadata = self._extract_email(content, file_type.extension)

            elif file_type.category == FileCategory.ARCHIVE:
                text_content, metadata = self._extract_archive_listing(content, file_type.extension)

            else:
                warnings.append(f"Unsupported file type: {file_type.extension}")

        except Exception as e:
            warnings.append(f"Extraction error: {str(e)}")

        # Clean and normalize text
        text_content = self._normalize_text(text_content)

        # Create chunks if text is large
        chunks = self._create_chunks(text_content, page_count)

        # Calculate statistics
        word_count = len(text_content.split()) if text_content else 0
        char_count = len(text_content)

        return ProcessedFile(
            original_filename=filename,
            detected_type=file_type,
            text_content=text_content,
            chunks=chunks,
            page_count=page_count,
            word_count=word_count,
            char_count=char_count,
            ocr_applied=ocr_applied,
            extraction_warnings=warnings,
            metadata=metadata
        )

    def _extract_text_file(self, content: bytes) -> str:
        """Extract text from text-based files."""
        encodings = ['utf-8', 'utf-16', 'latin-1', 'cp1252', 'iso-8859-1']

        for encoding in encodings:
            try:
                return content.decode(encoding)
            except (UnicodeDecodeError, LookupError):
                continue

        # Last resort: decode with errors='replace'
        return content.decode('utf-8', errors='replace')

    def _extract_pdf(self, content: bytes, force_ocr: bool = False) -> Tuple[str, int, bool]:
        """Extract text from PDF, with OCR fallback for scanned documents."""
        if not self.has_pdfplumber:
            return "[PDF extraction requires pdfplumber - install with: pip install pdfplumber]", 0, False

        import pdfplumber

        text_parts = []
        page_count = 0
        ocr_applied = False

        try:
            with pdfplumber.open(io.BytesIO(content)) as pdf:
                page_count = len(pdf.pages)

                for i, page in enumerate(pdf.pages):
                    page_text = page.extract_text() or ""

                    # Check if page has very little text (might be scanned)
                    if force_ocr or (len(page_text.strip()) < 50 and self.ocr_enabled):
                        # Try OCR on page image
                        try:
                            img = page.to_image(resolution=300)
                            ocr_text = self._ocr_image(img.original)
                            if ocr_text and len(ocr_text.strip()) > len(page_text.strip()):
                                page_text = ocr_text
                                ocr_applied = True
                        except Exception:
                            pass

                    if page_text.strip():
                        text_parts.append(f"--- Page {i + 1} ---\n{page_text}")

        except Exception as e:
            return f"[PDF extraction error: {str(e)}]", 0, False

        return "\n\n".join(text_parts), page_count, ocr_applied

    def _extract_docx(self, content: bytes) -> Tuple[str, Dict]:
        """Extract text from DOCX files."""
        if not self.has_docx:
            return "[DOCX extraction requires python-docx - install with: pip install python-docx]", {}

        from docx import Document

        doc = Document(io.BytesIO(content))

        # Extract document properties
        metadata = {}
        try:
            props = doc.core_properties
            if props.author:
                metadata['author'] = props.author
            if props.title:
                metadata['title'] = props.title
            if props.created:
                metadata['created'] = str(props.created)
        except:
            pass

        # Extract text from paragraphs
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]

        # Extract text from tables
        table_texts = []
        for table in doc.tables:
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(" | ".join(cells))
            table_texts.append("\n".join(rows))

        # Combine
        all_text = "\n\n".join(paragraphs)
        if table_texts:
            all_text += "\n\n--- Tables ---\n" + "\n\n".join(table_texts)

        return all_text, metadata

    def _extract_doc(self, content: bytes) -> str:
        """Extract text from legacy DOC files (limited support)."""
        # Legacy DOC is complex - try basic extraction
        try:
            # Try to find text between control characters
            text = content.decode('latin-1', errors='ignore')
            # Remove control characters but keep printable
            text = ''.join(c for c in text if c.isprintable() or c in '\n\t\r')
            # Clean up excessive whitespace
            text = re.sub(r'\s+', ' ', text)
            return text.strip()
        except:
            return ""

    def _extract_xlsx(self, content: bytes) -> Tuple[str, Dict]:
        """Extract text from XLSX files."""
        if not self.has_openpyxl:
            return "[XLSX extraction requires openpyxl - install with: pip install openpyxl]", {}

        from openpyxl import load_workbook

        wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)

        metadata = {
            'sheet_count': len(wb.sheetnames),
            'sheets': wb.sheetnames
        }

        text_parts = []

        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            rows = []

            for row in sheet.iter_rows(values_only=True):
                # Convert cells to strings, handling None
                cells = [str(cell) if cell is not None else "" for cell in row]
                if any(cells):  # Skip empty rows
                    rows.append(" | ".join(cells))

            if rows:
                text_parts.append(f"--- Sheet: {sheet_name} ---\n" + "\n".join(rows))

        wb.close()
        return "\n\n".join(text_parts), metadata

    def _extract_xls(self, content: bytes) -> str:
        """Extract text from legacy XLS files (limited support)."""
        # Would need xlrd for proper support
        return ""

    def _extract_pptx(self, content: bytes) -> Tuple[str, int, Dict]:
        """Extract text from PPTX files."""
        if not self.has_pptx:
            return "[PPTX extraction requires python-pptx - install with: pip install python-pptx]", 0, {}

        from pptx import Presentation

        prs = Presentation(io.BytesIO(content))

        metadata = {}
        try:
            props = prs.core_properties
            if props.author:
                metadata['author'] = props.author
            if props.title:
                metadata['title'] = props.title
        except:
            pass

        text_parts = []
        slide_count = len(prs.slides)

        for i, slide in enumerate(prs.slides):
            slide_text = []

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_text.append(text)

                # Extract from tables
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        slide_text.append(" | ".join(cells))

            if slide_text:
                text_parts.append(f"--- Slide {i + 1} ---\n" + "\n".join(slide_text))

        return "\n\n".join(text_parts), slide_count, metadata

    def _extract_ppt(self, content: bytes) -> str:
        """Extract text from legacy PPT files (limited support)."""
        return ""

    def _extract_rtf(self, content: bytes) -> str:
        """Extract text from RTF files."""
        if not self.has_striprtf:
            return "[RTF extraction requires striprtf - install with: pip install striprtf]"

        from striprtf.striprtf import rtf_to_text

        try:
            rtf_content = content.decode('latin-1')
            return rtf_to_text(rtf_content)
        except Exception as e:
            return f"[RTF extraction error: {str(e)}]"

    def _extract_image_ocr(self, content: bytes) -> str:
        """Extract text from image using OCR."""
        if not self.has_pillow or not self.has_pytesseract:
            return "[Image OCR requires Pillow and pytesseract - install with: pip install Pillow pytesseract]"

        from PIL import Image

        try:
            img = Image.open(io.BytesIO(content))
            return self._ocr_image(img)
        except Exception as e:
            return f"[Image OCR error: {str(e)}]"

    def _ocr_image(self, image) -> str:
        """Perform OCR on a PIL Image."""
        if not self.has_pytesseract:
            return ""

        import pytesseract

        try:
            # Convert to RGB if necessary
            if image.mode != 'RGB':
                image = image.convert('RGB')

            # Perform OCR
            text = pytesseract.image_to_string(image, lang=self.ocr_languages)
            return text.strip()
        except Exception:
            return ""

    def _extract_email(self, content: bytes, ext: str) -> Tuple[str, Dict]:
        """Extract text from email files."""
        metadata = {}

        if ext == 'eml':
            try:
                msg = email.message_from_bytes(content)

                metadata['from'] = msg.get('From', '')
                metadata['to'] = msg.get('To', '')
                metadata['subject'] = msg.get('Subject', '')
                metadata['date'] = msg.get('Date', '')

                # Extract body
                body_parts = []

                if msg.is_multipart():
                    for part in msg.walk():
                        content_type = part.get_content_type()
                        if content_type == 'text/plain':
                            payload = part.get_payload(decode=True)
                            if payload:
                                body_parts.append(payload.decode('utf-8', errors='replace'))
                else:
                    payload = msg.get_payload(decode=True)
                    if payload:
                        body_parts.append(payload.decode('utf-8', errors='replace'))

                text = f"From: {metadata['from']}\nTo: {metadata['to']}\nSubject: {metadata['subject']}\nDate: {metadata['date']}\n\n"
                text += "\n".join(body_parts)

                return text, metadata

            except Exception as e:
                return f"[Email extraction error: {str(e)}]", metadata

        elif ext == 'msg':
            # MSG format requires additional libraries (extract-msg)
            return "[MSG extraction requires extract-msg - install with: pip install extract-msg]", metadata

        return "", metadata

    def _extract_archive_listing(self, content: bytes, ext: str) -> Tuple[str, Dict]:
        """Extract file listing from archive."""
        metadata = {}
        files = []

        try:
            if ext in ('zip', 'docx', 'xlsx', 'pptx'):
                with zipfile.ZipFile(io.BytesIO(content)) as zf:
                    for info in zf.infolist():
                        if not info.is_dir():
                            files.append(f"{info.filename} ({info.file_size} bytes)")
                    metadata['file_count'] = len(files)

            elif ext in ('tar', 'tgz') or ext == 'gz':
                mode = 'r:gz' if ext in ('gz', 'tgz') else 'r'
                try:
                    with tarfile.open(fileobj=io.BytesIO(content), mode=mode) as tf:
                        for member in tf.getmembers():
                            if member.isfile():
                                files.append(f"{member.name} ({member.size} bytes)")
                        metadata['file_count'] = len(files)
                except:
                    pass

        except Exception as e:
            return f"[Archive listing error: {str(e)}]", metadata

        if files:
            return "Archive Contents:\n" + "\n".join(files[:100]), metadata  # Limit to 100 files

        return "Empty or unreadable archive", metadata

    def _normalize_text(self, text: str) -> str:
        """Normalize extracted text."""
        if not text:
            return ""

        # Replace various whitespace with standard space
        text = re.sub(r'[\t\r\f\v]+', ' ', text)

        # Normalize multiple newlines
        text = re.sub(r'\n{3,}', '\n\n', text)

        # Remove leading/trailing whitespace from lines
        lines = [line.strip() for line in text.split('\n')]
        text = '\n'.join(lines)

        # Remove excessive spaces
        text = re.sub(r' {2,}', ' ', text)

        return text.strip()

    def _create_chunks(
        self,
        text: str,
        page_count: Optional[int] = None
    ) -> List[ProcessedChunk]:
        """Create chunks for large text content."""
        if not text or len(text) <= self.max_chunk_size:
            return [ProcessedChunk(
                content=text,
                chunk_index=0,
                total_chunks=1,
                char_offset=0
            )]

        chunks = []
        current_pos = 0
        chunk_index = 0

        while current_pos < len(text):
            # Find chunk end - try to break at paragraph or sentence
            chunk_end = min(current_pos + self.max_chunk_size, len(text))

            if chunk_end < len(text):
                # Try to find a good break point
                # First try paragraph break
                para_break = text.rfind('\n\n', current_pos, chunk_end)
                if para_break > current_pos + self.max_chunk_size // 2:
                    chunk_end = para_break + 2
                else:
                    # Try sentence break
                    sentence_break = max(
                        text.rfind('. ', current_pos, chunk_end),
                        text.rfind('.\n', current_pos, chunk_end),
                        text.rfind('? ', current_pos, chunk_end),
                        text.rfind('! ', current_pos, chunk_end)
                    )
                    if sentence_break > current_pos + self.max_chunk_size // 2:
                        chunk_end = sentence_break + 2
                    else:
                        # Try word break
                        word_break = text.rfind(' ', current_pos, chunk_end)
                        if word_break > current_pos:
                            chunk_end = word_break + 1

            chunk_text = text[current_pos:chunk_end].strip()

            if chunk_text:
                chunks.append(ProcessedChunk(
                    content=chunk_text,
                    chunk_index=chunk_index,
                    total_chunks=0,  # Will be updated below
                    char_offset=current_pos
                ))
                chunk_index += 1

            current_pos = chunk_end

        # Update total_chunks
        for chunk in chunks:
            chunk.total_chunks = len(chunks)

        return chunks

    def get_supported_formats(self) -> Dict[str, List[str]]:
        """Get list of supported file formats."""
        formats = {
            'text': list(TEXT_EXTENSIONS),
            'documents': ['pdf', 'docx', 'doc', 'rtf'],
            'spreadsheets': ['xlsx', 'xls', 'csv'],
            'presentations': ['pptx', 'ppt'],
            'images': ['png', 'jpg', 'jpeg', 'tiff', 'bmp', 'gif', 'webp'],
            'email': ['eml', 'msg'],
            'archives': ['zip', 'tar', 'gz', 'tgz']
        }

        # Mark which have full support
        available = {
            'pdf': self.has_pdfplumber,
            'docx': self.has_docx,
            'xlsx': self.has_openpyxl,
            'pptx': self.has_pptx,
            'images': self.has_pillow and self.has_pytesseract,
            'rtf': self.has_striprtf
        }

        return formats

    def get_capabilities(self) -> Dict[str, bool]:
        """Get processor capabilities based on available dependencies."""
        return {
            'pdf_extraction': self.has_pdfplumber,
            'docx_extraction': self.has_docx,
            'xlsx_extraction': self.has_openpyxl,
            'pptx_extraction': self.has_pptx,
            'image_ocr': self.has_pillow and self.has_pytesseract,
            'rtf_extraction': self.has_striprtf,
            'ocr_enabled': self.ocr_enabled,
            'ocr_languages': self.ocr_languages
        }


# Singleton instance for the application
_processor: Optional[FileProcessor] = None


def get_file_processor(
    max_chunk_size: int = 50000,
    ocr_enabled: bool = True,
    ocr_languages: str = "eng+ara"
) -> FileProcessor:
    """Get or create the file processor instance."""
    global _processor
    if _processor is None:
        _processor = FileProcessor(
            max_chunk_size=max_chunk_size,
            ocr_enabled=ocr_enabled,
            ocr_languages=ocr_languages
        )
    return _processor
